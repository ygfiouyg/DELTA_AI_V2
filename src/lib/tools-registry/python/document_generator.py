"""
Tool: document_generator
Category: utility/docs
Package: python-docx, python-pptx, openpyxl, reportlab, fpdf2
Description: إنشاء مستندات Word, Excel, PowerPoint, PDF من بيانات.

Dependencies:
  - python-docx (pip install python-docx)
  - python-pptx (pip install python-pptx)
  - openpyxl (pip install openpyxl)
  - reportlab (pip install reportlab)
  - fpdf2 (pip install fpdf2)

Input:
  {
    "format": "docx" | "xlsx" | "pptx" | "pdf",
    "output_path": "/tmp/document.docx",
    "title": "Document Title",
    "content": { ... }  # format-specific
  }

Output:
  {"success": true, "file": "/tmp/document.docx", "size_kb": 25.3}
"""
import sys
import os
import json

for p in ["/usr/local/lib/python3.11/dist-packages", "/app/.venv/lib/python3.12/site-packages"]:
    if os.path.isdir(p) and p not in sys.path:
        sys.path.insert(0, p)


def generate(format: str, output_path: str, title: str = "", content: dict = None):
    content = content or {}
    if not output_path:
        return {"success": False, "error": "output_path required"}

    try:
        if format == "docx":
            return _generate_docx(output_path, title, content)
        elif format == "xlsx":
            return _generate_xlsx(output_path, title, content)
        elif format == "pptx":
            return _generate_pptx(output_path, title, content)
        elif format == "pdf":
            return _generate_pdf(output_path, title, content)
        else:
            return {"success": False, "error": f"unsupported format: {format}"}
    except Exception as e:
        return {"success": False, "error": f"generation failed: {str(e)[:200]}"}


def _generate_docx(output_path: str, title: str, content: dict):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor

    doc = Document()
    if title:
        h = doc.add_heading(title, level=0)
        h.alignment = 1  # center

    # Paragraphs
    for para in content.get("paragraphs", []):
        if isinstance(para, dict):
            text = para.get("text", "")
            style = para.get("style", "Normal")
            p = doc.add_paragraph(text, style=style)
        else:
            doc.add_paragraph(str(para))

    # Headings
    for h_data in content.get("headings", []):
        if isinstance(h_data, dict):
            doc.add_heading(h_data.get("text", ""), level=h_data.get("level", 1))
        else:
            doc.add_heading(str(h_data), level=1)

    # Table
    if content.get("table"):
        table_data = content["table"]
        if table_data and isinstance(table_data, list):
            rows = len(table_data)
            cols = len(table_data[0]) if table_data[0] else 0
            table = doc.add_table(rows=rows, cols=cols)
            table.style = "Table Grid"
            for r, row in enumerate(table_data):
                for c, cell in enumerate(row):
                    table.cell(r, c).text = str(cell)

    # Bullet list
    for item in content.get("bullets", []):
        doc.add_paragraph(str(item), style="List Bullet")

    doc.save(output_path)
    return {
        "success": True,
        "file": output_path,
        "size_kb": round(os.path.getsize(output_path) / 1024, 2),
        "format": "docx",
    }


def _generate_xlsx(output_path: str, title: str, content: dict):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.chart import BarChart, LineChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Sheet1"

    # Title row
    if title:
        ws["A1"] = title
        ws["A1"].font = Font(bold=True, size=14)
        ws.merge_cells("A1:E1")

    # Headers
    headers = content.get("headers", [])
    if headers:
        for col_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=2, column=col_idx, value=str(h))
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
            cell.alignment = Alignment(horizontal="center")

    # Data rows
    rows = content.get("rows", [])
    for r_idx, row in enumerate(rows, 3):
        for c_idx, val in enumerate(row, 1):
            ws.cell(row=r_idx, column=c_idx, value=val)

    # Auto-size columns (use cell.column_letter directly to avoid merged cell issues)
    from openpyxl.utils import get_column_letter
    for col_idx in range(1, len(headers) + 1):
        max_length = len(str(headers[col_idx - 1])) if col_idx <= len(headers) else 8
        for row in rows:
            if col_idx <= len(row):
                max_length = max(max_length, len(str(row[col_idx - 1])))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(50, max_length + 2)

    # Add chart if requested
    if content.get("chart") and rows:
        chart_type = content["chart"]
        if chart_type == "bar":
            chart = BarChart()
        elif chart_type == "line":
            chart = LineChart()
        else:
            chart = BarChart()
        data = Reference(ws, min_col=1, min_row=2, max_row=2 + len(rows), max_col=len(headers))
        chart.add_data(data, titles_from_data=True)
        chart.title = title or "Chart"
        ws.add_chart(chart, f"G2")

    wb.save(output_path)
    return {
        "success": True,
        "file": output_path,
        "size_kb": round(os.path.getsize(output_path) / 1024, 2),
        "format": "xlsx",
        "rows": len(rows),
    }


def _generate_pptx(output_path: str, title: str, content: dict):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title or "Presentation"
    if slide.placeholders[1]:
        slide.placeholders[1].text = content.get("subtitle", "")

    # Content slides
    for slide_data in content.get("slides", []):
        if isinstance(slide_data, dict):
            layout_idx = 1 if slide_data.get("bullets") else 5
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", "")
            if slide_data.get("bullets"):
                body = slide.placeholders[1]
                tf = body.text_frame
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

    prs.save(output_path)
    return {
        "success": True,
        "file": output_path,
        "size_kb": round(os.path.getsize(output_path) / 1024, 2),
        "format": "pptx",
        "slides": 1 + len(content.get("slides", [])),
    }


def _generate_pdf(output_path: str, title: str, content: dict):
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    doc = SimpleDocTemplate(output_path, pagesize=A4, topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    story = []

    if title:
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.2 * inch))

    for para in content.get("paragraphs", []):
        text = para if isinstance(para, str) else para.get("text", "")
        style_name = para.get("style", "Normal") if isinstance(para, dict) else "Normal"
        style = styles.get(style_name, styles["Normal"])
        story.append(Paragraph(text, style))
        story.append(Spacer(1, 0.1 * inch))

    if content.get("table"):
        table_data = content["table"]
        if table_data:
            # Wrap cells in Paragraph for proper rendering
            wrapped = []
            for row in table_data:
                wrapped.append([Paragraph(str(c), styles["Normal"]) for c in row])
            t = Table(wrapped)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4F46E5")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 12),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F3F4F6")),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(t)

    for bullet in content.get("bullets", []):
        story.append(Paragraph(f"• {bullet}", styles["Normal"]))

    doc.build(story)
    return {
        "success": True,
        "file": output_path,
        "size_kb": round(os.path.getsize(output_path) / 1024, 2),
        "format": "pdf",
    }



def _dispatch(args):
    return generate(args.get("format", ""), args.get("output_path", ""), args.get("title", ""), args.get("content", {}))


if __name__ == "__main__":

    # V.145: Support --args_file (called from Node.js registry)
    import sys as _sys
    if "--args_file" in _sys.argv:
        import json as _json
        _idx = _sys.argv.index("--args_file")
        with open(_sys.argv[_idx + 1]) as _f:
            _args = _json.load(_f)
        # Map args to function call based on script
        _result = _dispatch(_args)
        print(_json.dumps(_result, ensure_ascii=False, default=str))
        _sys.exit(0)

    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--format", required=True, choices=["docx", "xlsx", "pptx", "pdf"])
    parser.add_argument("--output_path", required=True)
    parser.add_argument("--title", default="")
    parser.add_argument("--content", default="{}", help="JSON string")
    args = parser.parse_args()
    content = json.loads(args.content)
    result = generate(args.format, args.output_path, args.title, content)
    print(json.dumps(result, ensure_ascii=False))
